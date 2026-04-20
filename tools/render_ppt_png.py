import argparse
import textwrap
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation


SLIDE_W = 1280
SLIDE_H = 720


def _load_font(preferred_name: str, size: int):
    try:
        return ImageFont.truetype(preferred_name, size)
    except Exception:
        return ImageFont.load_default()


def _wrap_lines(text: str, *, width_chars: int):
    lines = []
    for para in (text or "").splitlines():
        para = para.rstrip()
        if not para:
            lines.append("")
            continue
        wrapped = textwrap.wrap(para, width=width_chars, break_long_words=False, break_on_hyphens=False)
        lines.extend(wrapped if wrapped else [""])
    return lines


def _extract_slide_text(slide):
    title = ""
    if getattr(slide.shapes, "title", None) is not None:
        try:
            title = (slide.shapes.title.text or "").strip()
        except Exception:
            title = ""

    text_shapes = []
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        txt = (shape.text or "").strip()
        if not txt:
            continue
        text_shapes.append(txt)

    if not title and text_shapes:
        title = text_shapes[0].splitlines()[0].strip()

    body_parts = []
    for txt in text_shapes:
        if title and txt.strip() == title.strip():
            continue
        body_parts.append(txt)

    body = "\n\n".join(body_parts).strip()
    return title, body


def render_ppt_to_pngs(pptx_path: Path, out_dir: Path):
    prs = Presentation(str(pptx_path))
    out_dir.mkdir(parents=True, exist_ok=True)

    title_font = _load_font("arial.ttf", 44)
    body_font = _load_font("arial.ttf", 26)
    small_font = _load_font("arial.ttf", 18)

    for idx, slide in enumerate(prs.slides, start=1):
        title, body = _extract_slide_text(slide)
        img = Image.new("RGB", (SLIDE_W, SLIDE_H), (255, 255, 255))
        draw = ImageDraw.Draw(img)

        # Header bar
        draw.rectangle([0, 0, SLIDE_W, 110], fill=(15, 23, 42))
        draw.text((60, 30), title or f"Slide {idx}", font=title_font, fill=(255, 255, 255))

        # Body
        x = 80
        y = 150
        max_w = SLIDE_W - 160

        lines = _wrap_lines(body, width_chars=72)
        for line in lines:
            if y > SLIDE_H - 70:
                break
            if line == "":
                y += 14
                continue
            draw.text((x, y), line, font=body_font, fill=(30, 41, 59))
            y += 34

        # Footer
        footer = f"{idx}/{len(prs.slides)}"
        w = draw.textlength(footer, font=small_font)
        draw.text((SLIDE_W - 60 - w, SLIDE_H - 40), footer, font=small_font, fill=(100, 116, 139))

        out_path = out_dir / f"slide_{idx:02d}.png"
        img.save(out_path)


def main():
    parser = argparse.ArgumentParser(description="Render a PPTX to simple PNG previews.")
    parser.add_argument("--pptx", required=True, help="Path to .pptx file")
    parser.add_argument("--outdir", required=True, help="Output directory for PNGs")
    args = parser.parse_args()

    render_ppt_to_pngs(Path(args.pptx), Path(args.outdir))


if __name__ == "__main__":
    main()

