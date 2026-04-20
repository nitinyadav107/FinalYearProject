import argparse
from datetime import date, datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


PROJECT_TITLE = "GNN-Based Network Intrusion Detection System"
INSTITUTION = "NIET, Greater Noida"
COURSE = "B.Tech Computer Science"


def _set_run_style(run, *, bold=False, size=24, color=(20, 26, 38)):
    run.font.bold = bold
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor(*color)


def _add_title_slide(prs, *, title, subtitle_lines):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    subtitle = slide.placeholders[1].text_frame
    subtitle.clear()
    for i, line in enumerate(subtitle_lines):
        p = subtitle.paragraphs[0] if i == 0 else subtitle.add_paragraph()
        p.text = line
        p.level = 0
        for r in p.runs:
            _set_run_style(r, size=18, color=(71, 85, 105))
    return slide


def _add_bullets_slide(prs, *, title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        for r in p.runs:
            _set_run_style(r, size=20)
    return slide


def _add_two_column_slide(prs, *, title, left_title, left_bullets, right_title, right_bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title only
    slide.shapes.title.text = title

    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(4.3), Inches(4.8))
    left_tf = left_box.text_frame
    left_tf.clear()
    p0 = left_tf.paragraphs[0]
    run0 = p0.add_run()
    run0.text = left_title
    _set_run_style(run0, bold=True, size=20)
    for b in left_bullets:
        p = left_tf.add_paragraph()
        p.text = b
        p.level = 1
        for r in p.runs:
            _set_run_style(r, size=18)

    right_box = slide.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.3), Inches(4.8))
    right_tf = right_box.text_frame
    right_tf.clear()
    p1 = right_tf.paragraphs[0]
    run1 = p1.add_run()
    run1.text = right_title
    _set_run_style(run1, bold=True, size=20)
    for b in right_bullets:
        p = right_tf.add_paragraph()
        p.text = b
        p.level = 1
        for r in p.runs:
            _set_run_style(r, size=18)

    return slide


def create_presentation(*, students, out_path: Path):
    prs = Presentation()
    prs.core_properties.title = PROJECT_TITLE
    prs.core_properties.subject = "Final Year Project Presentation"
    prs.core_properties.created = datetime.now()

    team_lines = ["Group Members:"]
    for name, roll in students:
        team_lines.append(f"- {name} (Roll No: {roll})")
    team_lines += [COURSE, INSTITUTION, f"Date: {date.today().isoformat()}"]

    _add_title_slide(
        prs,
        title=PROJECT_TITLE,
        subtitle_lines=team_lines,
    )

    _add_bullets_slide(
        prs,
        title="Agenda",
        bullets=[
            "Problem statement & objectives",
            "Proposed approach using Graph Neural Networks (GNNs)",
            "System workflow and architecture",
            "Tech stack and models used",
            "Results (evaluation) and demo overview",
            "Conclusion and future work",
        ],
    )

    _add_bullets_slide(
        prs,
        title="Problem Statement",
        bullets=[
            "Traditional ML inspects network flows independently and may miss relationship-based attack patterns.",
            "Attacks like DDoS, reconnaissance, and lateral movement often appear as connected behavior across hosts.",
            "Goal: detect suspicious hosts by modeling traffic as a graph and learning from neighborhood context.",
        ],
    )

    _add_bullets_slide(
        prs,
        title="Objectives",
        bullets=[
            "Convert network traffic CSV flows into a host communication graph.",
            "Generate node-level features from traffic statistics (bytes, packets, duration, etc.).",
            "Train GNN models to classify nodes as benign vs. suspicious.",
            "Evaluate performance with standard metrics and visualize results.",
            "Provide an interactive dashboard for demonstration.",
        ],
    )

    _add_two_column_slide(
        prs,
        title="System Architecture (High Level)",
        left_title="Input & Preprocessing",
        left_bullets=[
            "Load CSV traffic data",
            "Normalize schema (src/dst IP, duration, label, etc.)",
            "Scale/encode features",
            "Build graph (nodes=hosts, edges=flows)",
        ],
        right_title="Learning & Output",
        right_bullets=[
            "Node feature aggregation",
            "Train GNN (GCN / GraphSAGE / GAT)",
            "Evaluate on test split",
            "Infer and rank suspicious hosts",
            "Show results in Streamlit dashboard",
        ],
    )

    _add_bullets_slide(
        prs,
        title="Workflow",
        bullets=[
            "Data collection / loading (real or synthetic)",
            "Schema normalization and label mapping",
            "Graph construction and node feature generation",
            "Train/test split and model training",
            "Evaluation (accuracy, precision, recall, F1, confusion matrix, ROC)",
            "Inference: suspicious host ranking",
        ],
    )

    _add_two_column_slide(
        prs,
        title="Models and Tech Stack",
        left_title="GNN Models",
        left_bullets=[
            "GCN (baseline graph convolution)",
            "GraphSAGE (neighborhood aggregation)",
            "GAT (attention-based aggregation)",
        ],
        right_title="Tools/Libraries",
        right_bullets=[
            "Python, PyTorch, PyTorch Geometric",
            "Pandas, NumPy, scikit-learn",
            "Matplotlib/Seaborn for plots",
            "Streamlit for dashboard",
        ],
    )

    _add_bullets_slide(
        prs,
        title="Results (Add Your Metrics)",
        bullets=[
            "Best model: <GCN / GraphSAGE / GAT>",
            "Accuracy: <..%>, Precision: <..%>, Recall: <..%>, F1-score: <..%>",
            "Confusion matrix summary: <TP/FP/TN/FN>",
            "Key observation: GNN captures neighbor context, improving detection of connected attacks.",
        ],
    )

    _add_bullets_slide(
        prs,
        title="Demo (Streamlit Dashboard)",
        bullets=[
            "Upload/select dataset and run preprocessing",
            "Train model or load trained weights",
            "View evaluation plots and metrics",
            "Inspect suspicious host ranking and graph insights",
        ],
    )

    _add_bullets_slide(
        prs,
        title="Conclusion & Future Work",
        bullets=[
            "Graph-based modeling improves intrusion detection by using relationships between hosts.",
            "GNNs learn both host features and neighborhood behavior for suspicious node prediction.",
            "Future work: larger real-world datasets, temporal graphs, explainability, and deployment hardening.",
            "Thank you!",
        ],
    )

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))


def _parse_students(values):
    students = []
    for v in values:
        if "|" in v:
            name, roll = [s.strip() for s in v.split("|", 1)]
        else:
            name, roll = v.strip(), "________"
        students.append((name or "________", roll or "________"))
    while len(students) < 3:
        students.append(("________", "________"))
    return students[:3]


def main():
    parser = argparse.ArgumentParser(description="Create a 10-slide final year project PPTX.")
    parser.add_argument(
        "--student",
        action="append",
        default=[],
        help='Student in "Name|RollNumber" format. Provide up to 3 times.',
    )
    parser.add_argument(
        "--out",
        default=str(Path.cwd() / "NIET_FYP_Presentation.pptx"),
        help="Output .pptx path",
    )
    args = parser.parse_args()

    students = _parse_students(args.student)
    create_presentation(students=students, out_path=Path(args.out))


if __name__ == "__main__":
    main()
