from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE_DIR = Path(__file__).resolve().parent
OUTPUT_PATH = BASE_DIR / "GNN_IDS_Final_Presentation.pptx"
METRICS_PATH = BASE_DIR / "metrics.json"
TRAINING_CURVE = BASE_DIR / "training_curves.png"
CONFUSION_MATRIX = BASE_DIR / "confusion_matrix.png"
ROC_CURVE = BASE_DIR / "roc_curve.png"

BG = RGBColor(247, 244, 235)
ACCENT = RGBColor(19, 78, 74)
ACCENT_2 = RGBColor(195, 80, 59)
TEXT = RGBColor(33, 37, 41)
MUTED = RGBColor(92, 99, 106)


def add_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.33), Inches(0.45))
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.8), Inches(11.8), Inches(0.9))
    paragraph = title_box.text_frame.paragraphs[0]
    run = paragraph.add_run()
    run.text = title
    run.font.name = "Aptos Display"
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = ACCENT

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.72), Inches(1.5), Inches(11.5), Inches(0.5))
        sub_paragraph = sub_box.text_frame.paragraphs[0]
        sub_run = sub_paragraph.add_run()
        sub_run.text = subtitle
        sub_run.font.name = "Aptos"
        sub_run.font.size = Pt(12)
        sub_run.font.color.rgb = MUTED


def add_bullets(slide, bullets: list[str], left: float = 0.9, top: float = 2.0, width: float = 11.0, height: float = 4.5) -> None:
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    text_frame = box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = 0
    text_frame.margin_right = 0

    for index, bullet in enumerate(bullets):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(21)
        paragraph.font.color.rgb = TEXT
        paragraph.space_after = Pt(10)


def add_two_column_bullets(slide, left_bullets: list[str], right_bullets: list[str]) -> None:
    add_bullets(slide, left_bullets, left=0.8, top=2.0, width=5.7, height=4.6)
    add_bullets(slide, right_bullets, left=6.8, top=2.0, width=5.5, height=4.6)


def add_image(slide, image_path: Path, left: float, top: float, width: float) -> None:
    if image_path.exists():
        slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width))


def add_footer(slide, text: str = "Final Year Project | GNN-Based Intrusion Detection") -> None:
    footer = slide.shapes.add_textbox(Inches(0.7), Inches(6.9), Inches(12.0), Inches(0.3))
    paragraph = footer.text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.RIGHT
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Aptos"
    run.font.size = Pt(10)
    run.font.color.rgb = MUTED


def metrics_text() -> list[str]:
    if not METRICS_PATH.exists():
        return [
            "Metrics file not found",
            "Run training first to populate evaluation values",
        ]

    metrics = json.loads(METRICS_PATH.read_text(encoding="utf-8"))
    return [
        f"Model used: {metrics.get('model', 'graphsage')}",
        f"Accuracy: {metrics.get('test_accuracy', 0):.3f}",
        f"Precision: {metrics.get('test_precision', 0):.3f}",
        f"Recall: {metrics.get('test_recall', 0):.3f}",
        f"F1 Score: {metrics.get('test_f1', 0):.3f}",
        f"ROC-AUC: {metrics.get('roc_auc') if metrics.get('roc_auc') is not None else 'N/A'}",
    ]


def build_presentation() -> Path:
    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    layout = presentation.slide_layouts[6]

    slides = []
    for _ in range(12):
        slide = presentation.slides.add_slide(layout)
        add_background(slide)
        slides.append(slide)

    add_title(slides[0], "Graph Neural Network Based Network Intrusion Detection System", "Fully completed project presentation")
    cover_box = slides[0].shapes.add_textbox(Inches(0.8), Inches(2.0), Inches(6.0), Inches(3.0))
    cover_tf = cover_box.text_frame
    for index, line in enumerate(
        [
            "Purpose: Automatically detect malicious network behavior",
            "Core Idea: Convert traffic into a graph and classify suspicious nodes",
            "Built With: Python, PyTorch Geometric, Streamlit, Matplotlib",
        ]
    ):
        paragraph = cover_tf.paragraphs[0] if index == 0 else cover_tf.add_paragraph()
        paragraph.text = line
        paragraph.font.name = "Aptos"
        paragraph.font.size = Pt(24 if index == 0 else 20)
        paragraph.font.color.rgb = TEXT
        paragraph.space_after = Pt(12)
    accent = slides[0].shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.0), Inches(1.8), Inches(4.2), Inches(3.0))
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT_2
    accent.line.fill.background()
    inner = slides[0].shapes.add_textbox(Inches(8.3), Inches(2.25), Inches(3.6), Inches(2.1))
    p = inner.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = "Nodes = IPs\nEdges = Communication\nOutput = Suspicious Hosts"
    r.font.name = "Aptos"
    r.font.size = Pt(22)
    r.font.bold = True
    r.font.color.rgb = RGBColor(255, 255, 255)
    add_footer(slides[0])

    slide_specs = [
        ("Problem Statement", "Why this project is needed", [
            "Large networks generate huge traffic and manual monitoring is difficult",
            "Traditional methods often inspect each traffic row independently",
            "Connected attacks like botnets and scans spread across multiple hosts",
            "Security teams need faster and more intelligent automatic detection",
        ]),
        ("Project Overview", "What the system does", [
            "Converts traffic logs into a communication graph",
            "Represents IP addresses or hosts as nodes",
            "Represents communication between hosts as edges",
            "Uses a Graph Neural Network to predict suspicious behavior",
        ]),
        ("Real Problem Solved", "Practical cybersecurity impact", [
            "Detects suspicious communication bursts and DDoS-like behavior",
            "Finds scanning and reconnaissance patterns",
            "Helps identify compromised or malicious hosts",
            "Reduces analyst workload by ranking risky nodes automatically",
        ]),
        ("Why Graph Neural Networks", "Why GNN is better here", [
            "Normal neural networks work on flat or sequential data",
            "Network traffic is relational and naturally graph-shaped",
            "GNN learns from both host features and neighbor behavior",
            "This makes it strong for multi-host and structural attack detection",
        ]),
        ("System Workflow", "End-to-end pipeline", [
            "Collect traffic data from CSV or monitoring tools",
            "Normalize fields like source IP, destination IP, bytes, packets, duration, and label",
            "Build graph edges from host-to-host communication",
            "Aggregate incoming and outgoing node features",
            "Train GCN, GraphSAGE, or GAT",
            "Evaluate metrics and run inference for suspicious node ranking",
        ]),
        ("Dataset Used", "Training and demonstration data", [
            "Synthetic dataset is used for guaranteed demo and testing",
            "Real CSV support is available for CICIDS-style traffic exports",
            "Logical fields required: src_ip, dst_ip, bytes, packets, duration, label",
            "Benign labels are mapped to 0 and attack labels are mapped to 1",
        ]),
        ("Technology Stack", "Tools and their roles", [
            "Python: complete project implementation",
            "PyTorch: tensor operations and model training",
            "PyTorch Geometric: graph neural network layers and graph data handling",
            "Pandas and scikit-learn: preprocessing, splits, metrics, scaling",
            "Matplotlib and Seaborn: charts and evaluation visuals",
            "Streamlit: dashboard for visual demo and inference",
        ]),
        ("Models and Their Roles", "What each model contributes", [
            "GCN: baseline graph convolution model for node classification",
            "GraphSAGE: scalable neighborhood aggregation for graph learning",
            "GAT: attention-based model that weighs important neighbors",
            "Training exports checkpoints, reports, plots, and ranked predictions",
        ]),
    ]

    for idx, (title, subtitle, bullets) in enumerate(slide_specs, start=1):
        add_title(slides[idx], title, subtitle)
        add_bullets(slides[idx], bullets)
        add_footer(slides[idx])

    add_title(slides[9], "Results from Tested Run", "Synthetic test execution results")
    add_two_column_bullets(slides[9], metrics_text(), [
        "Training run completed successfully in the tested environment",
        "Model checkpoint and reports were generated",
        "Inference produced ranked suspicious nodes",
        "The project is ready for demo through the Streamlit dashboard",
    ])
    add_image(slides[9], TRAINING_CURVE, left=7.1, top=3.3, width=5.4)
    add_footer(slides[9])

    add_title(slides[10], "Visual Outputs", "Generated project artifacts")
    add_image(slides[10], CONFUSION_MATRIX, left=0.8, top=1.9, width=4.0)
    add_image(slides[10], ROC_CURVE, left=4.7, top=1.9, width=4.0)
    add_image(slides[10], TRAINING_CURVE, left=8.5, top=1.9, width=4.0)
    add_bullets(slides[10], [
        "Confusion matrix shows classification quality",
        "ROC curve shows discrimination ability when both classes are present",
        "Training curves show how loss and F1 evolved during training",
    ], left=0.8, top=5.9, width=12.0, height=0.8)
    add_footer(slides[10])

    add_title(slides[11], "Conclusion and Future Scope", "What this project delivers")
    add_bullets(slides[11], [
        "The project provides a complete graph-based intrusion detection prototype",
        "It explains network attacks using structure, not only flat traffic rows",
        "It is useful for academic demonstration and can be extended for real deployments",
        "Future work: real-time capture, temporal GNNs, explainable AI, SOC integration",
    ])
    add_footer(slides[11])

    presentation.save(OUTPUT_PATH)
    return OUTPUT_PATH


if __name__ == "__main__":
    path = build_presentation()
    print(f"Saved presentation to {path}")
